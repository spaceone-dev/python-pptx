# encoding: utf-8

"""
Gherkin step implementations for chart features.
"""

from __future__ import absolute_import, print_function

from behave import given, then, when

from pptx import Presentation
from pptx.chart.axis import CategoryAxis, ValueAxis
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL_TYPE, MSO_THEME_COLOR

from .helpers import test_pptx


# given ===================================================

@given('a bar chart')
def given_a_bar_chart(context):
    prs = Presentation(test_pptx('cht-charts'))
    sld = prs.slides[0]
    graphic_frame = sld.shapes[0]
    context.chart = graphic_frame.chart


@given('a bar plot {having_or_not} data labels')
def given_a_bar_plot_having_or_not_data_labels(context, having_or_not):
    slide_idx = {
        'having':     0,
        'not having': 1,
    }[having_or_not]
    prs = Presentation(test_pptx('cht-plot-props'))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given('a bar plot having gap width of {width}')
def given_a_bar_plot_having_gap_width_of_width(context, width):
    slide_idx = {'no explicit value': 0, '300': 1}[width]
    prs = Presentation(test_pptx('cht-plot-props'))
    context.plot = prs.slides[slide_idx].shapes[0].chart.plots[0]


@given('a bar series having fill of {fill}')
def given_a_bar_series_having_fill_of_fill(context, fill):
    series_idx = {
        'Automatic': 0,
        'No Fill':   1,
        'Orange':    2,
        'Accent 1':  3,
    }[fill]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('a bar series having invert_if_negative of {setting}')
def given_a_bar_series_having_invert_if_negative_setting(context, setting):
    series_idx = {
        'no explicit setting': 0,
        'True':                1,
        'False':               2,
    }[setting]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('a bar series having {width} line')
def given_a_bar_series_having_width_line(context, width):
    series_idx = {
        'no':      0,
        '1 point': 1,
    }[width]
    prs = Presentation(test_pptx('cht-series-props'))
    plot = prs.slides[0].shapes[0].chart.plots[0]
    context.series = plot.series[series_idx]


@given('an axis having {major_or_minor} gridlines')
def given_an_axis_having_major_or_minor_gridlines(context, major_or_minor):
    prs = Presentation(test_pptx('cht-axis-props'))
    chart = prs.slides[0].shapes[0].chart
    context.axis = chart.value_axis


@given('an axis not having {major_or_minor} gridlines')
def given_an_axis_not_having_major_or_minor_gridlines(context, major_or_minor):
    prs = Presentation(test_pptx('cht-axis-props'))
    chart = prs.slides[0].shapes[0].chart
    context.axis = chart.category_axis


# when ====================================================

@when('I assign {value} to axis.has_{major_or_minor}_gridlines')
def when_I_assign_value_to_axis_has_major_or_minor_gridlines(
        context, value, major_or_minor):
    axis = context.axis
    propname = 'has_%s_gridlines' % major_or_minor
    new_value = {'True': True, 'False': False}[value]
    setattr(axis, propname, new_value)


@when('I assign {value} to plot.gap_width')
def when_I_assign_value_to_plot_gap_width(context, value):
    new_value = int(value)
    context.plot.gap_width = new_value


@when('I assign {value} to plot.has_data_labels')
def when_I_assign_value_to_plot_has_data_labels(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.plot.has_data_labels = new_value


@when('I assign {value} to series.invert_if_negative')
def when_I_assign_value_to_series_invert_if_negative(context, value):
    new_value = {
        'True':  True,
        'False': False,
    }[value]
    context.series.invert_if_negative = new_value


# then ====================================================

@then('axis.has_{major_or_minor}_gridlines is {value}')
def then_axis_has_major_or_minor_gridlines_is_expected_value(
        context, major_or_minor, value):
    axis = context.axis
    actual_value = {
        'major': axis.has_major_gridlines,
        'minor': axis.has_minor_gridlines,
    }[major_or_minor]
    expected_value = {'True': True, 'False': False}[value]
    assert actual_value is expected_value, 'got %s' % actual_value


@then('I can access the chart category axis')
def then_I_can_access_the_chart_category_axis(context):
    category_axis = context.chart.category_axis
    assert isinstance(category_axis, CategoryAxis)


@then('I can access the chart value axis')
def then_I_can_access_the_chart_value_axis(context):
    value_axis = context.chart.value_axis
    assert isinstance(value_axis, ValueAxis)


@then('the plot.has_data_labels property is {value}')
def then_the_plot_has_data_labels_property_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    assert context.plot.has_data_labels is expected_value


@then('the series fill RGB color is FF6600')
def then_the_series_fill_RGB_color_is_FF6600(context):
    fill = context.series.fill
    assert fill.fore_color.rgb == RGBColor(0xFF, 0x66, 0x00)


@then('the series fill theme color is Accent 1')
def then_the_series_fill_theme_color_is_Accent_1(context):
    fill = context.series.fill
    assert fill.fore_color.theme_color == MSO_THEME_COLOR.ACCENT_1


@then('the series has a fill type of {fill_type}')
def then_the_series_has_a_fill_type_of_type(context, fill_type):
    expected_fill_type = {
        'None':                     None,
        'MSO_FILL_TYPE.BACKGROUND': MSO_FILL_TYPE.BACKGROUND,
        'MSO_FILL_TYPE.SOLID':      MSO_FILL_TYPE.SOLID,
    }[fill_type]
    fill = context.series.fill
    assert fill.type == expected_fill_type


@then('the series has a line width of {width}')
def then_the_series_has_a_line_width_of_width(context, width):
    expected_width = int(width)
    line = context.series.line
    assert line.width == expected_width


@then('series.invert_if_negative is {value}')
def then_series_invert_if_negative_is_value(context, value):
    expected_value = {
        'True':  True,
        'False': False,
    }[value]
    series = context.series
    assert series.invert_if_negative is expected_value


@then('the value of plot.gap_width is {value}')
def then_the_value_of_plot_gap_width_is_value(context, value):
    expected_value = int(value)
    actual_gap_width = context.plot.gap_width
    assert actual_gap_width == expected_value, 'got %s' % actual_gap_width
